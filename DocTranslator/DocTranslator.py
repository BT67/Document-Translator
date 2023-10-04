import argparse
from pptx import Presentation
from googletrans import Translator
from httpcore import SyncHTTPProxy
from base64 import b64encode
from docx import Document
from pptx.util import Pt
from openpyxl import load_workbook

# Important:
# googletrans library must be version: googletrans==3.1.0a0 or higher
# docx must be installed using [pip install python-docx] not [pip install docx]

# def build_proxy_headers(username, password):
#     userpass = (username.encode("utf-8"), password.encode("utf-8"))
#     token = b64encode(b":".join(userpass))
#     return [(b"Proxy-Authorization", b"Basic " + token)]
#
#
# port = 8080
# proxy_url = (b"http", b"<proxy_host_ip>", port, b'')
# proxy_headers = build_proxy_headers("<username>", "<password>")
# proxy = {"https": SyncHTTPProxy(proxy_url=proxy_url, proxy_headers=proxy_headers)}

# translator = Translator(service_urls=['translate.googleapis.com'], proxies=proxy)

translator = Translator(service_urls=['translate.googleapis.com'])


def translate_spreadsheet(file, from_lang, to_lang):
    print(f"Translating spreadsheet: {file} from {from_lang} to {to_lang}")
    workbook = load_workbook(file)
    for worksheet in workbook.worksheets:
        for column in worksheet.columns:
            for cell in column:
                if cell.value is not None:
                    cell.value = translator.translate(cell.value, dest=to_lang).text
    # for worksheet in workbook.sheetnames:
    #     worksheet.title = translator.translate(worksheet.title, dest=to_lang).text
    filename = file[:file.rfind('.')] + "_" + to_lang + ".xlsx"
    workbook.save(filename)


def translate_word_document(file, from_lang, to_lang):
    print(f"Translating word document: {file} from {from_lang} to {to_lang}")
    document = Document(file)
    for paragraph in document.paragraphs:
        paragraph.text = translator.translate(paragraph.text, dest=to_lang).text
    filename = file[:file.rfind('.')] + "_" + to_lang + ".docx"
    document.save(filename)


def translate_text_file(file, from_lang, to_lang):
    print(f"Translating text file: {file} from {from_lang} to {to_lang}")


def translate_slideshow(file, from_lang, to_lang):
    print(f"Translating slideshow: {file} from {from_lang} to {to_lang} \n")
    presentation = Presentation(file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if len(shape.text) > 1:
                    shape.text = translator.translate(shape.text, dest=to_lang).text
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
    filename = file[:file.rfind('.')] + "_" + to_lang + ".pptx"
    presentation.save(filename)


def main():
    parser = argparse.ArgumentParser(description='Translate Documents')
    parser.add_argument('-i', '--input', help='Document to be translated')
    parser.add_argument('-f', '--from_lang', help='Document original language to be translated from (default Japanese)')
    parser.add_argument('-t', '--to_lang', help='Document target language to be translated to (default English)')
    args = parser.parse_args()

    if args.input is None:
        parser.print_help(file=None)
        return
    if args.from_lang is None:
        from_lang = 'ja'
    else:
        from_lang = args.from_lang
    if args.to_lang is None:
        to_lang = 'en'
    else:
        to_lang = args.to_lang

    file_ext = args.input[args.input.rfind('.'):]
    print(file_ext)
    match file_ext:
        case ".txt":
            translate_text_file(args.input, from_lang, to_lang)
        case ".docx":
            translate_word_document(args.input, from_lang, to_lang)
        case ".xlsx":
            translate_spreadsheet(args.input, from_lang, to_lang)
        case ".pptx":
            translate_slideshow(args.input, from_lang, to_lang)


if __name__ == "__main__":
    main()

from openpyxl.utils import get_column_letter
from pptx import Presentation
import PyPDF2 as pydf
from PyPDF2 import PdfReader, PdfWriter
from googletrans import Translator
from httpcore import SyncHTTPProxy
from base64 import b64encode
from docx import Document
from pptx.util import Pt
from openpyxl import load_workbook

# Important:
# googletrans library must be version: googletrans==3.1.0a0 or higher
# docx must be installed using [pip install python-docx] not [pip install docx]

# def build_proxy_headers(username, password):
#     userpass = (username.encode("utf-8"), password.encode("utf-8"))
#     token = b64encode(b":".join(userpass))
#     return [(b"Proxy-Authorization", b"Basic " + token)]
#
#
# port = 8080
# proxy_url = (b"http", b"<proxy_host_ip>", port, b'')
# proxy_headers = build_proxy_headers("<username>", "<password>")
# proxy = {"https": SyncHTTPProxy(proxy_url=proxy_url, proxy_headers=proxy_headers)}

# translator = Translator(service_urls=['translate.googleapis.com'], proxies=proxy)

translator = Translator(service_urls=['translate.googleapis.com'])


def case001():
    translate_slideshow("test.pptx", "ja", "en")


def case002():
    translate_word_document("test.docx", "ja", "en")


def case003():
    translate_spreadsheet("test.xlsx", "ja", "en")

def case004():
    translate_pdf("test.pdf", "ja", "en")


def translate_spreadsheet(file, from_lang, to_lang):
    print(f"Translating spreadsheet: {file} from {from_lang} to {to_lang}")
    workbook = load_workbook(file)
    for worksheet in workbook.worksheets:
        for column in worksheet.columns:
            for cell in column:
                if cell.value is not None:
                    cell.value = translator.translate(cell.value, dest=to_lang).text
    # for worksheet in workbook.sheetnames:
    #     worksheet.title = translator.translate(worksheet.title, dest=to_lang).text
    filename = file[:file.rfind('.')] + "_" + to_lang + ".xlsx"
    workbook.save(filename)


def translate_word_document(file, from_lang, to_lang):
    print(f"Translating word document: {file} from {from_lang} to {to_lang}")
    document = Document(file)
    for paragraph in document.paragraphs:
        paragraph.text = translator.translate(paragraph.text, dest=to_lang).text
    filename = file[:file.rfind('.')] + "_" + to_lang + ".docx"
    document.save(filename)


def translate_slideshow(file, from_lang, to_lang):
    print(f"Translating slideshow: {file} from {from_lang} to {to_lang} \n")
    presentation = Presentation(file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if len(shape.text) > 1:
                    shape.text = translator.translate(shape.text, dest=to_lang).text
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
    filename = file[:file.rfind('.')] + "_" + to_lang + ".pptx"
    presentation.save(filename)

def translate_pdf(file, from_lang, to_lang):
    print(f"Translating pdf document: {file} from {from_lang} to {to_lang}")
    pdf = PdfReader(file)
    writer = PdfWriter()
    for pageNum in range(0, len(pdf.pages) - 1):
        text = pdf.pages[pageNum].extract_text()
        if not text is None:
            text = translator.translate(text, dest=to_lang).text
            writer.add_page()
            writer.multi_cell(200, 10, txt=text, align="L")
    with open(file[:file.rfind('.')] + "_" + to_lang + ".pdf", 'wb') as out_file:
        writer.write(out_file)

def main():
    case001()
    case002()
    case003()
    case004()

if __name__ == "__main__":
    main()
